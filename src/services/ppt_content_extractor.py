"""
PowerPoint Content Extractor Service

Extracts text and content from PowerPoint presentations for LLM processing.
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional, Union
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PPT content extraction will be limited.")

@dataclass
class SlideContent:
    """Represents content from a single slide."""
    slide_number: int
    title: str
    text_content: List[str]
    notes: str
    layout_name: str
    shape_count: int

@dataclass
class PresentationContent:
    """Represents content from entire presentation."""
    file_path: str
    title: str
    total_slides: int
    slides: List[SlideContent]
    metadata: Dict

class PPTContentExtractor:
    """Service for extracting content from PowerPoint files."""

    def __init__(self):
        self.logger = logging.getLogger(__name__)
        if not PPTX_AVAILABLE:
            self.logger.warning("python-pptx not available. Limited functionality.")

    def extract_from_file(self, file_path: Union[str, Path]) -> Optional[PresentationContent]:
        """Extract content from a PowerPoint file."""
        if not PPTX_AVAILABLE:
            self.logger.error("Cannot extract content: python-pptx not available")
            return None

        try:
            file_path = Path(file_path)
            if not file_path.exists():
                self.logger.error(f"File not found: {file_path}")
                return None

            if file_path.suffix.lower() not in ['.pptx', '.ppt']:
                self.logger.error(f"Unsupported file format: {file_path.suffix}")
                return None

            self.logger.info(f"Extracting content from: {file_path}")

            # Handle .ppt files differently than .pptx
            if file_path.suffix.lower() == '.ppt':
                self.logger.warning(f"Legacy .ppt format detected: {file_path}")
                return self._extract_from_ppt_legacy(file_path)

            # Load presentation (for .pptx files)
            prs = Presentation(str(file_path))

            # Extract metadata
            metadata = self._extract_metadata(prs)

            # Extract slides
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_content = self._extract_slide_content(slide, i)
                slides.append(slide_content)

            presentation_title = metadata.get('title', file_path.stem)

            return PresentationContent(
                file_path=str(file_path),
                title=presentation_title,
                total_slides=len(slides),
                slides=slides,
                metadata=metadata
            )

        except Exception as e:
            self.logger.error(f"Failed to extract content from {file_path}: {e}")
            return None

    def _extract_from_ppt_legacy(self, file_path: Path) -> Optional[PresentationContent]:
        """Handle legacy .ppt files with fallback methods."""
        self.logger.info(f"Attempting legacy .ppt extraction: {file_path}")

        # First, try to load with python-pptx (sometimes works for newer .ppt files)
        try:
            prs = Presentation(str(file_path))
            self.logger.info("Successfully loaded .ppt file with python-pptx")

            # Extract metadata
            metadata = self._extract_metadata(prs)

            # Extract slides
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_content = self._extract_slide_content(slide, i)
                slides.append(slide_content)

            presentation_title = metadata.get('title', file_path.stem)

            return PresentationContent(
                file_path=str(file_path),
                title=presentation_title,
                total_slides=len(slides),
                slides=slides,
                metadata=metadata
            )

        except Exception as e:
            self.logger.warning(f"python-pptx failed for .ppt file: {e}")

        # Try conversion to .pptx if converter is available
        converted_file = self._try_convert_to_pptx(file_path)
        if converted_file:
            try:
                result = self.extract_from_file(converted_file)
                if result:
                    result.metadata['original_file'] = str(file_path)
                    result.metadata['converted_from'] = 'ppt'
                    self.logger.info("Successfully extracted content from converted .pptx")
                    return result
            except Exception as e:
                self.logger.warning(f"Failed to extract from converted file: {e}")
            finally:
                # Clean up converted file
                self._cleanup_converted_file(converted_file)

        # Final fallback: Create minimal content based on file info
        return self._create_minimal_ppt_content(file_path)

    def _try_convert_to_pptx(self, ppt_path: Path) -> Optional[Path]:
        """Try to convert .ppt to .pptx for better extraction."""
        try:
            from .ppt_converter import ppt_converter

            if not ppt_converter.is_conversion_available():
                self.logger.debug("Conversion not available, skipping")
                return None

            # Create temporary directory for converted file
            import tempfile
            temp_dir = Path(tempfile.gettempdir()) / "cuepilot_conversions"
            temp_dir.mkdir(exist_ok=True)

            converted_path = ppt_converter.convert_ppt_to_pptx(ppt_path, temp_dir)
            if converted_path and converted_path.exists():
                self.logger.info(f"Successfully converted .ppt to .pptx: {converted_path}")
                return converted_path

        except Exception as e:
            self.logger.debug(f"Conversion attempt failed: {e}")

        return None

    def _cleanup_converted_file(self, file_path: Path):
        """Clean up temporary converted file."""
        try:
            if file_path and file_path.exists():
                file_path.unlink()
                self.logger.debug(f"Cleaned up temporary file: {file_path}")
        except Exception as e:
            self.logger.debug(f"Failed to cleanup temporary file: {e}")

    def _create_minimal_ppt_content(self, file_path: Path) -> PresentationContent:
        """Create minimal content when direct .ppt extraction fails."""
        self.logger.info("Creating minimal content for .ppt file - using window-based tracking")

        # Try to get slide count from window detection if available
        slide_count = 1
        try:
            # Import here to avoid circular imports
            import sys
            backup_path = Path(__file__).parent.parent.parent / "backup_old_structure"
            if backup_path.exists():
                sys.path.insert(0, str(backup_path))
                from window_detector import PowerPointWindowDetector

                detector = PowerPointWindowDetector()
                window = detector.get_active_powerpoint_window()
                if window:
                    slide_info = detector.extract_slide_info_from_title(window.title)
                    if slide_info.get('total_slides'):
                        slide_count = slide_info['total_slides']
                        self.logger.info(f"Detected {slide_count} slides from window")
        except Exception as e:
            self.logger.debug(f"Could not get slide count from window: {e}")

        # Create placeholder slides
        slides = []
        for i in range(1, slide_count + 1):
            slides.append(SlideContent(
                slide_number=i,
                title=f"Slide {i}",
                text_content=[f"Content from slide {i} (extracted via window tracking)"],
                notes="",
                layout_name="Unknown Layout",
                shape_count=0
            ))

        metadata = {
            'title': file_path.stem,
            'slide_count': slide_count,
            'format': 'legacy_ppt',
            'extraction_method': 'window_based'
        }

        return PresentationContent(
            file_path=str(file_path),
            title=file_path.stem,
            total_slides=slide_count,
            slides=slides,
            metadata=metadata
        )

    def _extract_metadata(self, prs) -> Dict:
        """Extract presentation metadata."""
        metadata = {}
        try:
            core_props = prs.core_properties
            metadata.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'keywords': core_props.keywords or '',
                'category': core_props.category or '',
                'comments': core_props.comments or '',
                'slide_count': len(prs.slides)
            })
        except Exception as e:
            self.logger.warning(f"Could not extract metadata: {e}")
            metadata['slide_count'] = len(prs.slides)

        return metadata

    def _extract_slide_content(self, slide, slide_number: int) -> SlideContent:
        """Extract content from a single slide."""
        title = ""
        text_content = []
        notes = ""

        try:
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        text_content.append(shape.text.strip())
                    elif hasattr(shape, 'text_frame') and shape.text_frame:
                        # Handle title and content placeholders
                        if shape.text_frame.text.strip():
                            if not title:  # First text is often the title
                                title = shape.text_frame.text.strip()
                            else:
                                text_content.append(shape.text_frame.text.strip())
                    else:
                        text_content.append(shape.text.strip())

                # Extract table content
                elif hasattr(shape, 'table') and shape.table:
                    table_text = self._extract_table_content(shape.table)
                    if table_text:
                        text_content.append(table_text)

            # Extract slide notes
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text = self._extract_notes(slide.notes_slide)
                if notes_text:
                    notes = notes_text

            # Get layout name
            layout_name = ""
            if hasattr(slide, 'slide_layout') and slide.slide_layout:
                layout_name = slide.slide_layout.name or "Unknown Layout"

        except Exception as e:
            self.logger.warning(f"Error extracting content from slide {slide_number}: {e}")

        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            notes=notes,
            layout_name=layout_name,
            shape_count=len(slide.shapes)
        )

    def _extract_table_content(self, table) -> str:
        """Extract text from table shapes."""
        try:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    if cell.text.strip():
                        cells.append(cell.text.strip())
                if cells:
                    rows.append(" | ".join(cells))
            return "\n".join(rows) if rows else ""
        except Exception:
            return ""

    def _extract_notes(self, notes_slide) -> str:
        """Extract text from slide notes."""
        try:
            notes_text = []
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    notes_text.append(shape.text.strip())
            return "\n".join(notes_text)
        except Exception:
            return ""

    def get_slide_text(self, presentation: PresentationContent, slide_number: int) -> Optional[str]:
        """Get combined text content for a specific slide."""
        if slide_number < 1 or slide_number > len(presentation.slides):
            return None

        slide = presentation.slides[slide_number - 1]
        content_parts = []

        if slide.title:
            content_parts.append(f"Title: {slide.title}")

        if slide.text_content:
            content_parts.append("Content:")
            content_parts.extend(slide.text_content)

        if slide.notes:
            content_parts.append(f"Notes: {slide.notes}")

        return "\n".join(content_parts) if content_parts else ""

    def get_all_text(self, presentation: PresentationContent) -> str:
        """Get all text content from the presentation."""
        all_text = []
        all_text.append(f"Presentation: {presentation.title}")
        all_text.append(f"Total Slides: {presentation.total_slides}")
        all_text.append("=" * 50)

        for slide in presentation.slides:
            all_text.append(f"\nSlide {slide.slide_number}:")
            if slide.title:
                all_text.append(f"Title: {slide.title}")
            if slide.text_content:
                all_text.append("Content:")
                all_text.extend(slide.text_content)
            if slide.notes:
                all_text.append(f"Notes: {slide.notes}")
            all_text.append("-" * 30)

        return "\n".join(all_text)

# Service instance
ppt_extractor = PPTContentExtractor()