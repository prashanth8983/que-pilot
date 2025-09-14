"""
Content processing utilities for presentations.
Handles text extraction, OCR, and content analysis.
"""

import cv2
import numpy as np
import pytesseract
from PIL import Image
from typing import List, Dict, Optional
from pptx.slide import Slide


class ContentProcessor:
    """Handles content extraction and processing from presentation slides."""
    
    def __init__(self):
        self.ocr_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789.,!?;:()[]{}"\'-+=/\\@#$%^&*<>|~`_ '
    
    def extract_slide_text(self, slide: Slide) -> str:
        """Extract text content from a slide."""
        text_content = []
        
        try:
            for shape in slide.shapes:
                extracted_text = self._extract_shape_text(shape)
                if extracted_text:
                    text_content.append(extracted_text)
            
            # Extract notes if available
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                try:
                    notes_text = self._extract_notes_text(slide.notes_slide)
                    if notes_text:
                        text_content.append(f"Notes: {notes_text}")
                except Exception:
                    pass
            
            return '\n'.join(text_content)
        except Exception as e:
            print(f"Error extracting slide text: {str(e)}")
            return ""
    
    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape."""
        try:
            # Direct text access
            if hasattr(shape, 'text') and shape.text:
                return shape.text.strip()
            
            # Text frame access
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ""
                    for run in paragraph.runs:
                        if run.text:
                            paragraph_text += run.text
                    if paragraph_text.strip():
                        text_parts.append(paragraph_text.strip())
                return '\n'.join(text_parts)
            
            # Table content
            if hasattr(shape, 'table') and shape.table:
                table_text = []
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = self._extract_shape_text(cell)
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        table_text.append(' | '.join(row_text))
                return '\n'.join(table_text)
        except Exception:
            pass
        return ""
    
    def _extract_notes_text(self, notes_slide) -> str:
        """Extract text from slide notes."""
        try:
            notes_text = []
            for shape in notes_slide.shapes:
                shape_text = self._extract_shape_text(shape)
                if shape_text:
                    notes_text.append(shape_text)
            return '\n'.join(notes_text)
        except Exception:
            return ""
    
    def extract_text_with_ocr(self, slide_image: np.ndarray) -> str:
        """Extract text using OCR with image preprocessing."""
        if slide_image is None:
            return ""
        
        try:
            # Convert to PIL image
            if len(slide_image.shape) == 3 and slide_image.shape[2] == 3:
                pil_image = Image.fromarray(cv2.cvtColor(slide_image, cv2.COLOR_BGR2RGB))
            else:
                pil_image = Image.fromarray(slide_image)
            
            # Enhance image for better OCR
            enhanced_image = self._enhance_image_for_ocr(pil_image)
            
            # Try OCR with custom config first
            try:
                extracted_text = pytesseract.image_to_string(enhanced_image, config=self.ocr_config)
                if extracted_text.strip():
                    return extracted_text.strip()
            except Exception:
                pass
            
            # Fallback to default OCR
            extracted_text = pytesseract.image_to_string(enhanced_image)
            return extracted_text.strip()
        
        except Exception as e:
            print(f"OCR extraction failed: {str(e)}")
            return ""
    
    def _enhance_image_for_ocr(self, pil_image: Image.Image) -> Image.Image:
        """Enhance image quality for better OCR results."""
        try:
            img_array = np.array(pil_image)
            
            # Convert to grayscale if needed
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            # Apply enhancement techniques
            enhanced = cv2.convertScaleAbs(gray, alpha=1.2, beta=10)
            blurred = cv2.GaussianBlur(enhanced, (1, 1), 0)
            _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            return Image.fromarray(thresh)
        except Exception:
            return pil_image
    
    def detect_objects_in_slide(self, slide_image: np.ndarray) -> List[Dict]:
        """Detect objects/shapes in slide image."""
        if slide_image is None:
            return []
        
        objects_detected = []
        
        try:
            gray = cv2.cvtColor(slide_image, cv2.COLOR_BGR2GRAY)
            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            for i, contour in enumerate(contours):
                if cv2.contourArea(contour) > 100:
                    x, y, w, h = cv2.boundingRect(contour)
                    object_info = {
                        'id': i,
                        'type': 'detected_shape',
                        'bounding_box': (x, y, w, h),
                        'area': cv2.contourArea(contour)
                    }
                    objects_detected.append(object_info)
        except Exception as e:
            print(f"Object detection failed: {str(e)}")
        
        return objects_detected
    
    def slide_to_image(self, slide: Slide, presentation_width: int, presentation_height: int) -> np.ndarray:
        """Convert slide to image."""
        try:
            # Use presentation dimensions
            width = int(presentation_width.inches * 96)  # 96 DPI
            height = int(presentation_height.inches * 96)
            
            # Create white background
            img = np.ones((height, width, 3), dtype=np.uint8) * 255
            
            try:
                import io
                from PIL import ImageDraw, ImageFont
                
                pil_image = Image.fromarray(img)
                draw = ImageDraw.Draw(pil_image)
                
                try:
                    font = ImageFont.truetype("Arial.ttf", 24)
                except:
                    font = ImageFont.load_default()
                
                # Extract and render text shapes
                y_offset = 50
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        
                        try:
                            x = int(shape.left.inches * 96) if hasattr(shape, 'left') else 50
                            y = int(shape.top.inches * 96) if hasattr(shape, 'top') else y_offset
                            
                            x = max(10, min(x, width - 100))
                            y = max(10, min(y, height - 50))
                        except:
                            x = 50
                            y = y_offset
                        
                        draw.text((x, y), text, fill=(0, 0, 0), font=font)
                        y_offset += 40
                
                img = np.array(pil_image)
            except Exception:
                pass
            
            return img
        except Exception as e:
            print(f"Error in slide_to_image: {str(e)}")
            # Return basic white image as fallback
            return np.ones((600, 800, 3), dtype=np.uint8) * 255
