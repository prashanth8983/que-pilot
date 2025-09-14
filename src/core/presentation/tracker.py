import os
from typing import List, Optional, Dict
from pptx import Presentation
from pptx.slide import Slide
from PIL import Image, ImageDraw, ImageFont
import cv2
import numpy as np
import pytesseract
from .detector import PowerPointWindowDetector

class PresentationTracker:
    def __init__(self, ppt_path: str = None, auto_detect: bool = False):
        # Validate input parameters
        if ppt_path is not None and not isinstance(ppt_path, (str, os.PathLike)):
            raise TypeError(f"ppt_path must be a string or PathLike object, not {type(ppt_path)}")

        self.ppt_path = str(ppt_path) if ppt_path is not None else None
        self.presentation = None
        self.current_slide_index = 0
        self.slides = []
        self.total_slides = 0
        self.auto_detect = auto_detect
        self.window_detector = PowerPointWindowDetector() if auto_detect else None
        self.auto_sync_enabled = False
        self.is_ppt_file = False  # Flag to track if this is a .ppt file
        self.window_based_tracking = False  # Flag for window-based tracking mode

        if ppt_path:
            self.load_presentation()

    def _detect_ppt_slide_count(self) -> int:
        """Try to detect the actual slide count in a PPT file using various methods."""
        if not self.ppt_path or not os.path.exists(self.ppt_path):
            return 0

        try:
            # Method 1: Try using olefile to read the PowerPoint binary structure
            try:
                import olefile
                if olefile.isOleFile(self.ppt_path):
                    with olefile.OleFileIO(self.ppt_path) as ole:
                        # Look for slide-related streams in the OLE file
                        stream_names = ole.listdir()
                        slide_count = 0

                        for stream in stream_names:
                            stream_name = '/'.join(stream) if isinstance(stream, list) else str(stream)
                            # Count streams that typically contain slide data
                            if 'slide' in stream_name.lower() or 'Slide' in stream_name:
                                slide_count += 1

                        if slide_count > 0:
                            print(f"🔍 OLE file analysis found {slide_count} slide-related streams")
                            return slide_count

            except (ImportError, Exception) as e:
                print(f"🔍 OLE file analysis not available: {e}")

            # Method 2: Try using python-pptx with different approaches
            try:
                # Sometimes PPT files can be read if we ignore certain errors
                from pptx import Presentation
                from pptx.exc import PackageNotFoundError

                # Try opening with different error handling
                pres = Presentation(self.ppt_path)
                slide_count = len(pres.slides)
                if slide_count > 0:
                    print(f"🔍 Python-pptx successfully read {slide_count} slides")
                    return slide_count

            except Exception as e:
                print(f"🔍 Python-pptx analysis failed: {e}")

            # Method 3: File size-based estimation (improved)
            try:
                file_size = os.path.getsize(self.ppt_path)
                # Improved estimation based on typical PPT file sizes
                # Small presentations: ~50-100KB per slide
                # Medium presentations: ~30-80KB per slide (compressed)
                # Large presentations: ~20-60KB per slide (highly compressed)

                if file_size > 200000:  # 200KB+
                    # Use more aggressive estimation for larger files
                    # Assume 50KB per slide on average for better detection
                    estimated_slides = min(max(int(file_size / 50000), 2), 100)  # Cap between 2-100 slides
                    print(f"🔍 File size estimation suggests ~{estimated_slides} slides (file size: {file_size/1024/1024:.1f}MB)")
                    return estimated_slides
                elif file_size > 100000:  # 100KB+
                    # Smaller files, more conservative
                    estimated_slides = min(max(int(file_size / 80000), 2), 20)
                    print(f"🔍 File size estimation suggests ~{estimated_slides} slides (file size: {file_size/1024:.0f}KB)")
                    return estimated_slides

            except Exception as e:
                print(f"🔍 File size analysis failed: {e}")

            # Method 4: Look for PowerPoint processes and try to get info via AppleScript
            if hasattr(self, 'window_detector') and self.window_detector:
                try:
                    # Try macOS AppleScript to get slide count even if no window is active
                    if self.window_detector.system == "Darwin":
                        macos_info = self.window_detector.get_powerpoint_slide_info_macos()
                        if macos_info.get('total_slides') and macos_info['total_slides'] > 1:
                            print(f"🔍 macOS AppleScript detected {macos_info['total_slides']} slides")
                            return macos_info['total_slides']
                except Exception as e:
                    print(f"🔍 AppleScript analysis failed: {e}")

            return 0

        except Exception as e:
            print(f"🔍 All slide count detection methods failed: {e}")
            return 0

    def load_presentation(self):
        if self.ppt_path is None:
            raise ValueError("No presentation path specified")

        if not isinstance(self.ppt_path, (str, os.PathLike)):
            raise TypeError(f"ppt_path must be a string or PathLike object, not {type(self.ppt_path)}")

        if not os.path.exists(self.ppt_path):
            raise FileNotFoundError(f"PowerPoint file not found: {self.ppt_path}")

        # Check file format
        file_extension = os.path.splitext(self.ppt_path)[1].lower()
        if file_extension not in ['.pptx', '.ppt']:
            raise Exception(f"Unsupported file format: {file_extension}")

        if file_extension == '.ppt':
            # Handle .ppt files - try PPTX loading first, fall back to window-based tracking
            self.is_ppt_file = True
            print(f"Attempting to load .ppt file: {os.path.basename(self.ppt_path)}")

            # Try to load as PPTX first (some .ppt files can be read by python-pptx)
            try:
                self.presentation = Presentation(self.ppt_path)
                self.slides = list(self.presentation.slides)
                self.total_slides = len(self.slides)
                self.window_based_tracking = False  # Direct file access works
                print(f"📊 Successfully loaded .ppt file as PPTX with {self.total_slides} slides")
                return
            except Exception as pptx_error:
                print(f"Could not load .ppt as PPTX: {pptx_error}")
                # Fall back to window-based tracking
                self.window_based_tracking = True
                self.presentation = None
                self.slides = []

            # If auto_detect is available, try to get slide info from window
            if self.auto_detect and self.window_detector:
                window_info = self.window_detector.get_current_slide_info()
                if window_info and window_info.get('total_slides'):
                    self.total_slides = window_info['total_slides']
                    if window_info.get('current_slide'):
                        self.current_slide_index = window_info['current_slide'] - 1
                    print(f"Loaded .ppt file with window-based tracking: {self.total_slides} slides detected")
                    return
                else:
                    # Try to detect actual slide count before falling back to default
                    actual_slide_count = self._detect_ppt_slide_count()
                    if actual_slide_count > 1:
                        self.total_slides = actual_slide_count
                        print(f"📊 Detected {self.total_slides} slides in .ppt file using file analysis")
                    else:
                        # Default values for .ppt files when all detection fails
                        self.total_slides = 1
                        print(f"Loaded .ppt file: Window-based tracking enabled (slide count will update when PowerPoint is detected)")
                    return
            else:
                # No auto-detect available, try to get slide count from file analysis
                actual_slide_count = self._detect_ppt_slide_count()
                if actual_slide_count > 1:
                    self.total_slides = actual_slide_count
                    print(f"📊 Detected {self.total_slides} slides in .ppt file using file analysis")
                else:
                    # Default fallback when all methods fail
                    self.total_slides = 1
                    print(f"Loaded .ppt file: Limited functionality (enable auto_detect for full tracking)")
                return

        # Handle .pptx files normally
        try:
            self.presentation = Presentation(self.ppt_path)
            self.slides = list(self.presentation.slides)
            self.total_slides = len(self.slides)
            print(f"📊 Loaded .pptx presentation with {self.total_slides} slides")
        except Exception as e:
            raise Exception(f"Error loading presentation: {str(e)}")

    def get_current_slide_number(self) -> int:
        return self.current_slide_index + 1

    def get_total_slides(self) -> int:
        # For .ppt files, try to get updated count from window if available
        if self.window_based_tracking and self.window_detector:
            window_info = self.window_detector.get_current_slide_info()
            if window_info and window_info.get('total_slides'):
                self.total_slides = window_info['total_slides']
        return self.total_slides

    def next_slide(self) -> bool:
        if self.current_slide_index < self.total_slides - 1:
            self.current_slide_index += 1
            return True
        return False

    def previous_slide(self) -> bool:
        if self.current_slide_index > 0:
            self.current_slide_index -= 1
            return True
        return False

    def go_to_slide(self, slide_number: int) -> bool:
        if 1 <= slide_number <= self.total_slides:
            self.current_slide_index = slide_number - 1
            return True
        return False

    def get_current_slide(self) -> Slide:
        return self.slides[self.current_slide_index]

    def extract_slide_text(self, slide_index: Optional[int] = None) -> str:
        """Extract text content from slide with improved handling"""
        if slide_index is None:
            slide_index = self.current_slide_index

        if not (0 <= slide_index < self.total_slides):
            return ""

        # For .ppt files, we can't extract text directly from file
        if self.window_based_tracking:
            return f"Text extraction from .ppt files requires OCR. Current slide: {slide_index + 1}"

        # For .pptx files, extract text normally
        if not self.slides or slide_index >= len(self.slides):
            return ""

        slide = self.slides[slide_index]
        text_content = []

        try:
            for shape in slide.shapes:
                extracted_text = self._extract_shape_text(shape)
                if extracted_text:
                    text_content.append(extracted_text)

            # Also check slide notes if available
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                try:
                    notes_text = self._extract_notes_text(slide.notes_slide)
                    if notes_text:
                        text_content.append(f"Notes: {notes_text}")
                except Exception:
                    pass  # Notes extraction failed, continue without notes

            return '\n'.join(text_content)

        except Exception as e:
            print(f"Error extracting slide text: {str(e)}")
            return ""

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape with comprehensive handling"""
        try:
            # Direct text access
            if hasattr(shape, 'text') and shape.text:
                return shape.text.strip()

            # Check if shape has text frame
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ""
                    for run in paragraph.runs:
                        if run.text:
                            paragraph_text += run.text
                    if paragraph_text.strip():
                        text_parts.append(paragraph_text.strip())
                return '\n'.join(text_parts)

            # Check for table content
            if hasattr(shape, 'table') and shape.table:
                table_text = []
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = self._extract_shape_text(cell)
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        table_text.append(' | '.join(row_text))
                return '\n'.join(table_text)

        except Exception:
            pass  # Shape text extraction failed

        return ""

    def _extract_notes_text(self, notes_slide) -> str:
        """Extract text from slide notes"""
        try:
            notes_text = []
            for shape in notes_slide.shapes:
                shape_text = self._extract_shape_text(shape)
                if shape_text:
                    notes_text.append(shape_text)
            return '\n'.join(notes_text)
        except Exception:
            return ""

    def get_slide_image(self, slide_index: Optional[int] = None) -> Optional[np.ndarray]:
        if slide_index is None:
            slide_index = self.current_slide_index

        if not (0 <= slide_index < self.total_slides):
            return None

        # For window-based tracking without slides data, we can't generate slide images
        if self.window_based_tracking or not self.slides:
            if self.is_ppt_file and not self.slides:
                print(f"Cannot generate slide image for slide {slide_index + 1}: .ppt file could not be loaded directly and requires window-based tracking")
            else:
                print(f"Cannot generate slide image for slide {slide_index + 1}: window-based tracking mode active")
            return None

        try:
            slide = self.slides[slide_index]
            slide_image = self._slide_to_image(slide)
            return slide_image
        except Exception as e:
            print(f"Error converting slide to image: {str(e)}")
            return None

    def _slide_to_image(self, slide: Slide) -> np.ndarray:
        """Convert slide to image using a more comprehensive approach"""
        try:
            # Use the presentation's slide dimensions
            width = int(self.presentation.slide_width.inches * 96)  # 96 DPI for better quality
            height = int(self.presentation.slide_height.inches * 96)

            # Create white background
            img = np.ones((height, width, 3), dtype=np.uint8) * 255

            # Try to extract shapes and render them (simplified approach)
            # This is a basic implementation - for full rendering, you'd need more complex logic
            try:
                import io
                from PIL import Image, ImageDraw, ImageFont

                # Create PIL image for drawing
                pil_image = Image.fromarray(img)
                draw = ImageDraw.Draw(pil_image)

                # Try to get a default font
                try:
                    font = ImageFont.truetype("Arial.ttf", 24)
                except:
                    font = ImageFont.load_default()

                # Extract text shapes and their positions
                y_offset = 50
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Simple text rendering at approximate positions
                        text = shape.text.strip()

                        # Try to get shape position if available
                        try:
                            x = int(shape.left.inches * 96) if hasattr(shape, 'left') else 50
                            y = int(shape.top.inches * 96) if hasattr(shape, 'top') else y_offset

                            # Ensure coordinates are within bounds
                            x = max(10, min(x, width - 100))
                            y = max(10, min(y, height - 50))
                        except:
                            x = 50
                            y = y_offset

                        # Draw text
                        draw.text((x, y), text, fill=(0, 0, 0), font=font)
                        y_offset += 40

                # Convert back to numpy array
                img = np.array(pil_image)

            except Exception as text_render_error:
                # If text rendering fails, still return the white background
                pass

            return img

        except Exception as e:
            print(f"Error in _slide_to_image: {str(e)}")
            # Return a basic white image as fallback
            width = 800
            height = 600
            return np.ones((height, width, 3), dtype=np.uint8) * 255

    def extract_text_with_ocr(self, slide_index: Optional[int] = None) -> str:
        """Extract text using OCR with improved image preprocessing"""
        slide_image = self.get_slide_image(slide_index)
        if slide_image is None:
            return ""

        try:
            # Convert to PIL image with proper color handling
            if len(slide_image.shape) == 3 and slide_image.shape[2] == 3:
                # BGR to RGB conversion for OpenCV images
                pil_image = Image.fromarray(cv2.cvtColor(slide_image, cv2.COLOR_BGR2RGB))
            else:
                pil_image = Image.fromarray(slide_image)

            # Enhance image for better OCR results
            enhanced_image = self._enhance_image_for_ocr(pil_image)

            # Configure Tesseract for better results
            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789.,!?;:()[]{}"\'-+=/\\@#$%^&*<>|~`_ '

            try:
                # Try with custom config first
                extracted_text = pytesseract.image_to_string(enhanced_image, config=custom_config)
                if extracted_text.strip():
                    return extracted_text.strip()
            except Exception:
                pass  # Fall back to default config

            # Default OCR as fallback
            extracted_text = pytesseract.image_to_string(enhanced_image)
            return extracted_text.strip()

        except Exception as e:
            print(f"OCR extraction failed: {str(e)}")
            # If OCR fails completely, try to extract text directly from shapes
            direct_text = self.extract_slide_text(slide_index)
            if direct_text:
                return f"Direct extraction (OCR failed): {direct_text}"
            return ""

    def _enhance_image_for_ocr(self, pil_image: Image.Image) -> Image.Image:
        """Enhance image quality for better OCR results"""
        try:
            # Convert to numpy array for OpenCV processing
            img_array = np.array(pil_image)

            # Convert to grayscale if needed
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array

            # Apply image enhancement techniques
            # 1. Increase contrast
            enhanced = cv2.convertScaleAbs(gray, alpha=1.2, beta=10)

            # 2. Apply Gaussian blur to reduce noise
            blurred = cv2.GaussianBlur(enhanced, (1, 1), 0)

            # 3. Apply threshold to get better text contrast
            _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

            # Convert back to PIL Image
            return Image.fromarray(thresh)

        except Exception:
            # If enhancement fails, return original image
            return pil_image

    def detect_objects_in_slide(self, slide_index: Optional[int] = None) -> List[dict]:
        slide_image = self.get_slide_image(slide_index)
        if slide_image is None:
            return []

        objects_detected = []

        try:
            gray = cv2.cvtColor(slide_image, cv2.COLOR_BGR2GRAY)

            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

            for i, contour in enumerate(contours):
                if cv2.contourArea(contour) > 100:
                    x, y, w, h = cv2.boundingRect(contour)

                    object_info = {
                        'id': i,
                        'type': 'detected_shape',
                        'bounding_box': (x, y, w, h),
                        'area': cv2.contourArea(contour)
                    }
                    objects_detected.append(object_info)

        except Exception as e:
            print(f"Object detection failed: {str(e)}")

        return objects_detected

    def get_slide_info(self, slide_index: Optional[int] = None) -> dict:
        if slide_index is None:
            slide_index = self.current_slide_index

        slide_text = self.extract_slide_text(slide_index)
        ocr_text = self.extract_text_with_ocr(slide_index)
        detected_objects = self.detect_objects_in_slide(slide_index)

        return {
            'slide_number': slide_index + 1,
            'total_slides': self.total_slides,
            'text_content': slide_text,
            'ocr_text': ocr_text,
            'detected_objects': detected_objects,
            'object_count': len(detected_objects)
        }

    def search_text_in_slides(self, search_term: str) -> List[int]:
        matching_slides = []

        for i in range(self.total_slides):
            slide_text = self.extract_slide_text(i).lower()
            ocr_text = self.extract_text_with_ocr(i).lower()

            if (search_term.lower() in slide_text or
                search_term.lower() in ocr_text):
                matching_slides.append(i + 1)

        return matching_slides

    def auto_detect_presentation_file(self) -> Optional[str]:
        """Auto-detect presentation file with enhanced search"""
        if not self.window_detector:
            return None

        window_info = self.get_window_info()
        if window_info and window_info['slide_info'].get('presentation_name'):
            ppt_name = window_info['slide_info']['presentation_name']

            common_paths = [
                os.path.expanduser("~/Desktop"),
                os.path.expanduser("~/Documents"),
                os.path.expanduser("~/Downloads"),
                os.getcwd()
            ]

            for path in common_paths:
                if os.path.exists(path):
                    for file in os.listdir(path):
                        if file.lower().endswith(('.ppt', '.pptx')):
                            if ppt_name.lower() in file.lower() or file.lower() in ppt_name.lower():
                                full_path = os.path.join(path, file)
                                if not hasattr(self, '_last_detected_file') or self._last_detected_file != full_path:
                                    print(f"🎯 Auto-detected presentation file: {full_path}")
                                    self._last_detected_file = full_path
                                return full_path
        return None

    def get_window_info(self) -> Optional[Dict]:
        """Get comprehensive window information"""
        if not self.window_detector:
            return None

        window = self.window_detector.get_active_powerpoint_window()
        if window:
            try:
                slide_info = self.window_detector.extract_slide_info_from_title(window.title)
            except Exception as e:
                print(f"Error extracting slide info from title '{window.title}': {e}")
                slide_info = {'current_slide': 1, 'total_slides': 0, 'presentation_name': None, 'mode': 'error'}
            return {
                'window_title': window.title,
                'app_name': window.app_name,
                'slide_info': slide_info,
                'window_position': window.position,
                'window_size': window.size
            }
        return None

    def enable_auto_sync(self):
        if not self.window_detector:
            print("Auto-sync requires auto_detect=True during initialization")
            return False

        self.auto_sync_enabled = True
        return True

    def disable_auto_sync(self):
        self.auto_sync_enabled = False

    def sync_with_powerpoint_window(self) -> bool:
        if not self.window_detector:
            return False

        slide_info = self.window_detector.get_current_slide_info()
        if not slide_info or not slide_info.get('current_slide'):
            return False

        detected_slide = slide_info['current_slide']

        # For .ppt files, update total slides from window info
        if self.window_based_tracking and slide_info.get('total_slides'):
            self.total_slides = slide_info['total_slides']

        if 1 <= detected_slide <= self.total_slides:
            if detected_slide - 1 != self.current_slide_index:
                self.current_slide_index = detected_slide - 1
                print(f"Auto-synced to slide {detected_slide}")
                return True

        return False

    def get_window_info(self) -> Optional[Dict]:
        if not self.window_detector:
            return None

        window = self.window_detector.get_active_powerpoint_window()
        if window:
            try:
                slide_info = self.window_detector.extract_slide_info_from_title(window.title)
            except Exception as e:
                print(f"Error extracting slide info from title '{window.title}': {e}")
                slide_info = {'current_slide': 1, 'total_slides': 0, 'presentation_name': None, 'mode': 'error'}
            return {
                'window_title': window.title,
                'app_name': window.app_name,
                'slide_info': slide_info,
                'window_position': window.position,
                'window_size': window.size
            }
        return None

    def start_auto_monitoring(self, callback=None, interval: float = 1.0):
        if not self.window_detector:
            print("Auto-monitoring requires auto_detect=True during initialization")
            return

        def sync_callback(window, slide_info):
            if self.auto_sync_enabled and self.presentation:
                current_slide = slide_info.get('current_slide')
                if current_slide and 1 <= current_slide <= self.total_slides:
                    if current_slide - 1 != self.current_slide_index:
                        self.current_slide_index = current_slide - 1
                        print(f"🔄 Auto-synced to slide {current_slide}")

            if callback:
                callback(self, window, slide_info)

        self.window_detector.monitor_powerpoint_window(sync_callback, interval)

    def auto_detect_presentation_file(self) -> Optional[str]:
        if not self.window_detector:
            return None

        window_info = self.get_window_info()
        if window_info and window_info['slide_info'].get('presentation_name'):
            ppt_name = window_info['slide_info']['presentation_name']

            common_paths = [
                os.path.expanduser("~/Desktop"),
                os.path.expanduser("~/Documents"),
                os.path.expanduser("~/Downloads"),
                os.getcwd()
            ]

            for path in common_paths:
                if os.path.exists(path):
                    for file in os.listdir(path):
                        if file.lower().endswith(('.ppt', '.pptx')):
                            if ppt_name.lower() in file.lower() or file.lower() in ppt_name.lower():
                                full_path = os.path.join(path, file)
                                # Only print once when we first detect it
                                if not hasattr(self, '_last_detected_file') or self._last_detected_file != full_path:
                                    print(f"Auto-detected presentation file: {full_path}")
                                    self._last_detected_file = full_path
                                return full_path

        return None

    def auto_load_presentation(self) -> bool:
        detected_file = self.auto_detect_presentation_file()
        if detected_file and os.path.exists(detected_file):
            file_extension = os.path.splitext(detected_file)[1].lower()

            try:
                self.ppt_path = detected_file
                self.load_presentation()
                self.enable_auto_sync()

                if file_extension == '.ppt':
                    print(f"Loaded .ppt file with window-based tracking: {os.path.basename(detected_file)}")
                else:
                    print(f"📊 Loaded .pptx presentation: {os.path.basename(detected_file)}")

                return True
            except Exception as e:
                print(f"Error loading auto-detected presentation: {e}")
                return False
        return False

    def extract_text_with_ocr(self, slide_index: Optional[int] = None) -> str:
        """Extract text using OCR with improved image preprocessing"""
        slide_image = self.get_slide_image(slide_index)
        if slide_image is None:
            return ""

        try:
            # Convert to PIL image with proper color handling
            if len(slide_image.shape) == 3 and slide_image.shape[2] == 3:
                # BGR to RGB conversion for OpenCV images
                pil_image = Image.fromarray(cv2.cvtColor(slide_image, cv2.COLOR_BGR2RGB))
            else:
                pil_image = Image.fromarray(slide_image)

            # Enhance image for better OCR results
            enhanced_image = self._enhance_image_for_ocr(pil_image)

            # Configure Tesseract for better results
            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789.,!?;:()[]{}"\'-+=/\\@#$%^&*<>|~`_ '

            try:
                # Try with custom config first
                extracted_text = pytesseract.image_to_string(enhanced_image, config=custom_config)
                if extracted_text.strip():
                    return extracted_text.strip()
            except Exception:
                pass  # Fall back to default config

            # Default OCR as fallback
            extracted_text = pytesseract.image_to_string(enhanced_image)
            return extracted_text.strip()

        except Exception as e:
            print(f"OCR extraction failed: {str(e)}")
            # If OCR fails completely, try to extract text directly from shapes
            direct_text = self.extract_slide_text(slide_index)
            if direct_text:
                return f"Direct extraction (OCR failed): {direct_text}"
            return ""

    def search_text_in_slides(self, search_term: str) -> List[int]:
        """Search for text across all slides using both direct text extraction and OCR"""
        matching_slides = []

        for i in range(self.total_slides):
            slide_text = self.extract_slide_text(i).lower()
            ocr_text = self.extract_text_with_ocr(i).lower()

            if (search_term.lower() in slide_text or
                search_term.lower() in ocr_text):
                matching_slides.append(i + 1)

        return matching_slides

    def auto_detect_presentation_file(self) -> Optional[str]:
        """Auto-detect presentation file with enhanced search"""
        if not self.window_detector:
            return None

        window_info = self.get_window_info()
        if window_info and window_info['slide_info'].get('presentation_name'):
            ppt_name = window_info['slide_info']['presentation_name']

            common_paths = [
                os.path.expanduser("~/Desktop"),
                os.path.expanduser("~/Documents"),
                os.path.expanduser("~/Downloads"),
                os.getcwd()
            ]

            for path in common_paths:
                if os.path.exists(path):
                    for file in os.listdir(path):
                        if file.lower().endswith(('.ppt', '.pptx')):
                            if ppt_name.lower() in file.lower() or file.lower() in ppt_name.lower():
                                full_path = os.path.join(path, file)
                                if not hasattr(self, '_last_detected_file') or self._last_detected_file != full_path:
                                    print(f"🎯 Auto-detected presentation file: {full_path}")
                                    self._last_detected_file = full_path
                                return full_path
        return None

    def get_window_info(self) -> Optional[Dict]:
        """Get comprehensive window information"""
        if not self.window_detector:
            return None

        window = self.window_detector.get_active_powerpoint_window()
        if window:
            try:
                slide_info = self.window_detector.extract_slide_info_from_title(window.title)
            except Exception as e:
                print(f"Error extracting slide info from title '{window.title}': {e}")
                slide_info = {'current_slide': 1, 'total_slides': 0, 'presentation_name': None, 'mode': 'error'}
            return {
                'window_title': window.title,
                'app_name': window.app_name,
                'slide_info': slide_info,
                'window_position': window.position,
                'window_size': window.size
            }
        return None